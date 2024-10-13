from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def pdf_to_pptx(pdf_path, pptx_path):
    # PDFを画像に変換
    images = convert_from_path(pdf_path)

    # 新しいPPTXプレゼンテーションを作成
    prs = Presentation()

    for image in images:
        # 新しいスライドを追加
        slide_layout = prs.slide_layouts[5]  # レイアウト「5」は空のスライドです
        slide = prs.slides.add_slide(slide_layout)

        # スライドに画像を追加
        left = top = Inches(0)
        pic = slide.shapes.add_picture(image.filename, left, top, height=prs.slide_height)

    # PPTXファイルを保存
    prs.save(pptx_path)

# 使用例
pdf_path = '0715_mori.pdf'
pptx_path = 'sub.pptx'
pdf_to_pptx(pdf_path, pptx_path)
